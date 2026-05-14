"""Distill every file currently in _scrape_dump/downloads/ into Jerry's
knowledge file at jerry_gpt_knowledge/09_roadmap_documents.md.

Standalone — does NOT re-scrape. Use whenever you've manually added or
re-pulled documents and just want to refresh the Jerry-readable markdown.

Usage:
    cd salestoolkit
    source .scrape_venv/bin/activate
    python scripts/distill_downloads.py
"""
from __future__ import annotations

from datetime import datetime, timezone
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
DOWNLOADS_DIR = ROOT / "_scrape_dump" / "downloads"
OUTPUT_FILE = ROOT / "jerry_gpt_knowledge" / "09_roadmap_documents.md"

# Per-file text cap (chars). Bumped from 6K → 12K because Release Notes
# are dense — capability lists, fixed-bug enumerations, supported-vehicle
# tables. Truncating too early loses signal.
PER_DOC_CHARS = 12000


def extract_docx(path: Path) -> str:
    from docx import Document
    doc = Document(path)
    parts: list[str] = []
    for p in doc.paragraphs:
        t = p.text.strip()
        if t:
            parts.append(t)
        if sum(len(x) for x in parts) > PER_DOC_CHARS:
            break
    for tbl in doc.tables:
        for row in tbl.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
        if sum(len(x) for x in parts) > PER_DOC_CHARS:
            break
    return "\n\n".join(parts)[:PER_DOC_CHARS]


def extract_xlsx(path: Path) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(path, data_only=True, read_only=True)
    parts: list[str] = []
    for sn in wb.sheetnames:
        parts.append(f"### Sheet: {sn}")
        ws = wb[sn]
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            cells = ["" if v is None else str(v).strip() for v in row]
            if any(cells):
                parts.append(" | ".join(cells))
            if i > 300:
                parts.append("[truncated at 300 rows]")
                break
        if sum(len(x) for x in parts) > PER_DOC_CHARS:
            break
    return "\n".join(parts)[:PER_DOC_CHARS]


def extract_pdf(path: Path) -> str:
    try:
        import pdfplumber
    except ImportError:
        return "[pdfplumber not installed]"
    parts: list[str] = []
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages):
            t = page.extract_text() or ""
            if t.strip():
                parts.append(f"-- Page {i + 1} --\n{t.strip()}")
            if sum(len(x) for x in parts) > PER_DOC_CHARS:
                break
    return "\n\n".join(parts)[:PER_DOC_CHARS]


def extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        return "[python-pptx not installed]"
    prs = Presentation(path)
    parts: list[str] = []
    for i, slide in enumerate(prs.slides):
        slide_lines = [f"-- Slide {i + 1} --"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(r.text for r in para.runs).strip()
                    if t:
                        slide_lines.append(t)
        if len(slide_lines) > 1:
            parts.append("\n".join(slide_lines))
        if sum(len(x) for x in parts) > PER_DOC_CHARS:
            break
    return "\n\n".join(parts)[:PER_DOC_CHARS]


EXTRACTORS = {
    ".docx": extract_docx,
    ".xlsx": extract_xlsx,
    ".xlsm": extract_xlsx,
    ".pdf": extract_pdf,
    ".pptx": extract_pptx,
    ".txt": lambda p: p.read_text(encoding="utf-8", errors="ignore")[:PER_DOC_CHARS],
    ".md":  lambda p: p.read_text(encoding="utf-8", errors="ignore")[:PER_DOC_CHARS],
}


def main():
    if not DOWNLOADS_DIR.exists():
        print(f"ERROR: {DOWNLOADS_DIR} does not exist. Run the scraper first.")
        return 1

    files = sorted(
        [p for p in DOWNLOADS_DIR.iterdir()
         if p.is_file() and not p.name.startswith(".")],
        key=lambda p: p.name,
    )
    if not files:
        print(f"ERROR: no files in {DOWNLOADS_DIR}")
        return 1

    print(f"Distilling {len(files)} files from {DOWNLOADS_DIR}")
    extracted: list[dict] = []
    for f in files:
        ext = f.suffix.lower()
        size_kb = f.stat().st_size // 1024
        extractor = EXTRACTORS.get(ext)
        if not extractor:
            print(f"  skip {f.name} ({ext} unsupported)")
            extracted.append({"path": f, "text": f"[unsupported file type: {ext}]", "size_kb": size_kb})
            continue
        print(f"  extract {f.name} ({size_kb:,} KB)")
        try:
            text = extractor(f)
            extracted.append({"path": f, "text": text, "size_kb": size_kb})
            print(f"    {len(text):,} chars extracted")
        except Exception as e:
            print(f"    FAILED: {type(e).__name__}: {e}")
            extracted.append({"path": f, "text": f"[extraction failed: {type(e).__name__}: {e}]", "size_kb": size_kb})

    # Build markdown
    scraped_at = datetime.now(timezone.utc).isoformat(timespec="seconds")
    lines: list[str] = []
    lines.append("# Streamax Internal Version Portal — Downloaded Documents")
    lines.append("")
    lines.append(f"_Distilled at `{scraped_at}` from `{DOWNLOADS_DIR.relative_to(ROOT)}`._")
    lines.append("")
    lines.append(
        f"This file contains text extracted from {len(extracted)} documents downloaded "
        "from the internal Streamax version portal (路线规划 / 文档中心 / 版本日志). "
        "Release notes, feature lists, configuration manuals, and operation guides are "
        "all parsed; per-document text is capped at "
        f"~{PER_DOC_CHARS:,} characters. When a user asks about a specific spec, "
        "release note, supported feature, or operational procedure, ground the answer in "
        "this file's extracts. Reference documents by filename when relevant."
    )
    lines.append("")
    lines.append("---")
    lines.append("")
    lines.append("## Document index")
    lines.append("")
    lines.append("| File | Type | Size |")
    lines.append("| --- | --- | --- |")
    for e in extracted:
        p = e["path"]
        lines.append(f"| `{p.name}` | {p.suffix.lstrip('.').upper()} | {e['size_kb']:,} KB |")
    lines.append("")
    lines.append("---")
    lines.append("")

    for e in extracted:
        p = e["path"]
        lines.append(f"## File: `{p.name}`")
        lines.append("")
        text = e["text"].strip()
        if text:
            lines.append("```")
            lines.append(text)
            lines.append("```")
        else:
            lines.append("_(no text extracted)_")
        lines.append("")
        lines.append("---")
        lines.append("")

    OUTPUT_FILE.parent.mkdir(parents=True, exist_ok=True)
    md = "\n".join(lines)
    OUTPUT_FILE.write_text(md, encoding="utf-8")
    print(f"\nWrote {len(md):,} chars → {OUTPUT_FILE.relative_to(ROOT)}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
