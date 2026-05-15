"""Streamax internal roadmap-portal scraper for Jerry GPT.

Runs locally on your Mac (which is on the Streamax intranet) and pulls the
content from http://10.20.51.20:5173 — version portal sidebar with 8
sections: 路线规划, 版本履历表, 版本日志, 功能清单, 文档中心, 功能说明,
MDVR 用户操作配置, 常用链接.

Now also: finds every downloadable document linked from any of those pages,
fetches it with the authenticated session, and distills the contents (PDF,
Word, Excel, PowerPoint) into a second markdown file Jerry reads.

Usage:
    one-time setup
        python3 -m venv .scrape_venv
        source .scrape_venv/bin/activate
        pip install playwright pdfplumber python-docx openpyxl python-pptx
        playwright install chromium

    every time you want fresh content:
        source .scrape_venv/bin/activate
        python scrape_roadmap.py

Output:
    _scrape_dump/                              raw HTML + screenshots + API JSON
    _scrape_dump/downloads/                    every file downloaded from the portal
                                               (gitignored — local inspection only)
    jerry_gpt_knowledge/08_roadmap_portal.md   distilled page text Jerry reads
    jerry_gpt_knowledge/09_roadmap_documents.md  distilled document text Jerry reads
"""
from __future__ import annotations

import json
import re
import sys
from datetime import datetime, timezone
from pathlib import Path
from urllib.parse import unquote, urlparse

try:
    from playwright.sync_api import sync_playwright, TimeoutError as PWTimeout
except ImportError:
    print("ERROR: playwright is not installed.\n")
    print("Run these once:")
    print("    python3 -m venv .scrape_venv")
    print("    source .scrape_venv/bin/activate")
    print("    pip install playwright pdfplumber python-docx openpyxl python-pptx")
    print("    playwright install chromium")
    print("\nThen re-run: python scrape_roadmap.py")
    sys.exit(1)

# Optional document-extraction libraries. Each is tried/imported separately
# so a missing one only disables that file type, not the whole pipeline.
try:
    import pdfplumber
except ImportError:
    pdfplumber = None
try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None
try:
    from openpyxl import load_workbook
except ImportError:
    load_workbook = None
try:
    from pptx import Presentation as PptxPresentation
except ImportError:
    PptxPresentation = None


# --- Config -----------------------------------------------------------------

BASE_URL = "http://10.20.51.20:5173"
LOGIN_USER = "admin"
LOGIN_PASS = "admin"

SIDEBAR_ITEMS = [
    ("01_roadmap",            "路线规划",         "Roadmap (release timeline by month, version cards)"),
    ("02_version_history",    "版本履历表",       "Version history table"),
    ("03_version_log",        "版本日志",         "Version log / change notes"),
    ("04_feature_list",       "功能清单",         "Feature list / capability inventory"),
    ("05_docs_center",        "文档中心",         "Documentation center"),
    ("06_feature_description","功能说明",         "Feature descriptions / specs"),
    ("07_mdvr_config",        "MDVR 用户操作配置", "MDVR user operation configuration"),
    ("08_common_links",       "常用链接",         "Common links / external resources"),
]

POST_LOAD_MS = 2000

# Download caps so a runaway portal doesn't fill the disk
MAX_FILE_BYTES = 50 * 1024 * 1024     # skip any single file larger than 50 MB
MAX_TOTAL_FILES = 200                  # cap total downloads per run

# Per-document extraction cap so the distilled markdown stays manageable.
# Anthropic's context window is huge but Jerry's effective working memory
# benefits from focused content over raw dumps.
PER_DOC_TEXT_CHARS = 6000

DOC_EXTENSIONS = {
    ".pdf", ".docx", ".doc", ".xlsx", ".xls", ".pptx", ".ppt",
    ".zip", ".csv", ".txt", ".md",
}
DOWNLOAD_HINTS_IN_URL = ("/download", "/file/", "/attachment", "/files/")

ROOT = Path(__file__).resolve().parent
DUMP_DIR = ROOT / "_scrape_dump"
DOWNLOADS_DIR = DUMP_DIR / "downloads"
KNOWLEDGE_DIR = ROOT / "jerry_gpt_knowledge"
PAGES_OUTPUT_FILE = KNOWLEDGE_DIR / "08_roadmap_portal.md"
DOCS_OUTPUT_FILE = KNOWLEDGE_DIR / "09_roadmap_documents.md"


# --- Login flow -------------------------------------------------------------

def _try_login(page) -> bool:
    print(f"[scrape] navigating to {BASE_URL}")
    try:
        page.goto(BASE_URL, wait_until="networkidle", timeout=15000)
    except PWTimeout:
        page.wait_for_load_state("domcontentloaded", timeout=10000)
    page.wait_for_timeout(1500)

    if not page.locator("input[type='password']").count():
        print("[scrape] no password field visible — already logged in?")
        return True

    user_selectors = [
        'input[type="text"]',
        'input[name="username"]',
        'input[name="account"]',
        'input[placeholder*="账号"]',
        'input[placeholder*="用户"]',
        'input[placeholder*="user"]',
    ]
    pwd_selectors = ['input[type="password"]']
    submit_selectors = [
        'button[type="submit"]',
        'button:has-text("登录")',
        'button:has-text("登 录")',
        'button:has-text("Login")',
        'button:has-text("Sign in")',
        '.login-btn',
        'button.el-button--primary',
        'button.ant-btn-primary',
    ]

    filled_user = filled_pwd = False
    for sel in user_selectors:
        try:
            page.fill(sel, LOGIN_USER, timeout=1500)
            filled_user = True
            print(f"[scrape] filled username via '{sel}'")
            break
        except Exception:
            continue
    for sel in pwd_selectors:
        try:
            page.fill(sel, LOGIN_PASS, timeout=1500)
            filled_pwd = True
            print(f"[scrape] filled password via '{sel}'")
            break
        except Exception:
            continue

    if not (filled_user and filled_pwd):
        return False

    for sel in submit_selectors:
        try:
            page.click(sel, timeout=1500)
            print(f"[scrape] clicked submit via '{sel}'")
            try:
                page.wait_for_load_state("networkidle", timeout=10000)
            except PWTimeout:
                pass
            page.wait_for_timeout(2000)
            return True
        except Exception:
            continue
    try:
        page.locator('input[type="password"]').press("Enter")
        page.wait_for_load_state("networkidle", timeout=10000)
        page.wait_for_timeout(2000)
        return True
    except Exception:
        pass
    return False


# --- Sidebar walk + download discovery -------------------------------------

def _capture_section(page, slug: str, label: str, dump_dir: Path) -> dict:
    print(f"[scrape] section: {label}")
    try:
        item = page.get_by_text(label, exact=True).first
        item.click(timeout=5000)
        try:
            page.wait_for_load_state("networkidle", timeout=8000)
        except PWTimeout:
            pass
        page.wait_for_timeout(POST_LOAD_MS)

        captured_url = page.url
        html = page.content()
        visible_text = page.evaluate("() => document.body.innerText")

        (dump_dir / f"{slug}.html").write_text(html, encoding="utf-8")
        (dump_dir / f"{slug}.txt").write_text(visible_text, encoding="utf-8")
        page.screenshot(path=str(dump_dir / f"{slug}.png"), full_page=True)

        # Extract download links from this page
        download_links = _extract_download_urls_from_page(page)
        print(f"[scrape]   {len(visible_text):,} chars, {len(download_links)} download links")

        return {
            "label": label,
            "captured_url": captured_url,
            "visible_text": visible_text,
            "download_links": download_links,
            "ok": True,
        }
    except Exception as exc:
        print(f"[scrape]   FAILED ({type(exc).__name__}: {exc})")
        try:
            page.screenshot(path=str(dump_dir / f"{slug}_error.png"), full_page=True)
        except Exception:
            pass
        return {"label": label, "ok": False, "error": f"{type(exc).__name__}: {exc}"}


def _extract_download_urls_from_page(page) -> list[dict]:
    """Find every <a href> on the current page that looks like a download.

    Most portals don't expose downloads as <a href> — they use JS click
    handlers that trigger a Blob download. This function still runs as a
    cheap first pass, but the heavy lifting happens in
    `_harvest_docs_center()` which actively clicks cloud-download icons.
    """
    js = """
    (extensions, hints) => {
        const out = [];
        const seen = new Set();
        for (const a of document.querySelectorAll('a[href]')) {
            const href = a.href;
            const text = (a.innerText || a.title || '').trim();
            if (!href || href.startsWith('javascript:') || href.startsWith('#')) continue;
            const lower = href.toLowerCase();
            const looksLikeDoc = extensions.some(e => lower.endsWith(e)
                || lower.includes(e + '?'));
            const looksLikeDownload = hints.some(h => lower.includes(h));
            if ((looksLikeDoc || looksLikeDownload) && !seen.has(href)) {
                seen.add(href);
                out.push({url: href, text: text});
            }
        }
        return out;
    }
    """
    try:
        return page.evaluate(js, list(DOC_EXTENSIONS), list(DOWNLOAD_HINTS_IN_URL))
    except Exception:
        return []


def _harvest_docs_center(page, downloads_dir: Path, captured_downloads: list[dict]) -> int:
    """Specialized walk of 文档中心 — verbose logging so we can SEE failures.

    On the portal, each document row has a cloud-download icon (tooltip
    "下载原始 DOCX" or similar). The icon is gated behind:
        1. A category tab (Release Note / 功能说明 / 产品资料)
        2. A collapsed model row (e.g., "通用 1个版本")
        3. The cloud icon itself, at the far-right of the document row

    This function logs everything it tries — categories found, expansions
    attempted, every download candidate's tag/title/position, every click
    result. After a failed run, the console output tells us exactly which
    step needs adjustment, and `_scrape_dump/docs_center_*.png/html` give
    us the visual + DOM evidence to fix it.

    Downloads are caught by the context-level `download` listener already
    set up — that listener auto-saves them. We also capture popup events
    in case the portal opens downloads in a new tab instead.
    """
    print("[scrape]   ========== DOCS-CENTER HARVEST ==========")
    initial_count = len(captured_downloads)

    # Save the initial state so we can inspect what the script actually saw
    try:
        page.screenshot(path=str(DUMP_DIR / "docs_center_00_initial.png"), full_page=True)
        (DUMP_DIR / "docs_center_00_initial.html").write_text(page.content(), encoding="utf-8")
    except Exception:
        pass

    # Listen for popups too — some portals open download URLs in a new tab
    popup_urls: list[str] = []

    def _on_popup(p):
        try:
            popup_urls.append(p.url)
            print(f"[scrape]     ⊕ popup opened: {p.url[:100]}")
            p.close()
        except Exception:
            pass

    page.on("popup", _on_popup)

    # ---- 1. Find category tabs ----
    categories_text = ["Release Note", "功能说明", "产品资料", "技术文档", "用户手册"]
    category_locators: list[tuple[str, "object | None"]] = []
    for cat_text in categories_text:
        for sel in [
            f'text=/{cat_text}.*篇/',
            f'div:has-text("{cat_text}"):has-text("篇")',
            f'[class*="card"]:has-text("{cat_text}")',
        ]:
            try:
                loc = page.locator(sel).first
                if loc.count() > 0:
                    category_locators.append((cat_text, loc))
                    print(f"[scrape]   • category found '{cat_text}' via {sel!r}")
                    break
            except Exception:
                continue
    if not category_locators:
        print("[scrape]   ⚠ no category cards found — will try with whatever's visible")
        category_locators = [("(default view)", None)]

    # ---- 2. Walk each category ----
    total_click_attempts = 0
    for cat_idx, (cat_label, cat_loc) in enumerate(category_locators):
        print(f"\n[scrape]   ----- category {cat_idx + 1}/{len(category_locators)}: {cat_label} -----")
        if cat_loc is not None:
            try:
                cat_loc.scroll_into_view_if_needed(timeout=2000)
                cat_loc.click(timeout=3000)
                print(f"[scrape]     ✓ category click ok")
                page.wait_for_timeout(1500)
            except Exception as e:
                print(f"[scrape]     ✗ category click FAILED: {type(e).__name__}: {e}")
                continue

        # Save state after category click
        try:
            safe = re.sub(r"[^a-zA-Z0-9]+", "_", cat_label)
            page.screenshot(path=str(DUMP_DIR / f"docs_center_{cat_idx:02d}_{safe}.png"), full_page=True)
        except Exception:
            pass

        # ---- 2a. Expand all collapsed rows ----
        expand_selectors = [
            '[aria-expanded="false"]',
            '.el-collapse-item__header:not(.is-active)',
            'div[role="button"][aria-expanded="false"]',
            '[class*="collapse-header"]',
            # Generic: any clickable row that looks like a model header
            'div:has-text("个版本"):not(:has(div:has-text("个版本")))',
        ]
        expand_total = 0
        for sel in expand_selectors:
            try:
                rows = page.locator(sel).all()
                if rows:
                    print(f"[scrape]     expanding {len(rows)} rows via {sel!r}")
                    for row in rows:
                        try:
                            row.click(timeout=1000)
                            expand_total += 1
                            page.wait_for_timeout(150)
                        except Exception:
                            pass
            except Exception:
                continue
        if expand_total:
            print(f"[scrape]     {expand_total} expansion clicks total")
        page.wait_for_timeout(800)

        # ---- 2b. Find download candidates ----
        candidate_selectors = [
            '[title*="下载"]',
            '[title*="download"]',
            '[aria-label*="下载"]',
            '[aria-label*="download"]',
        ]
        # Collect (element, title, position) tuples; dedupe by position
        seen_positions = set()
        candidates: list[dict] = []
        for sel in candidate_selectors:
            try:
                els = page.locator(sel).all()
            except Exception:
                continue
            if els:
                print(f"[scrape]     {sel!r}: {len(els)} matches")
            for el in els:
                try:
                    box = el.bounding_box()
                    if not box or box["width"] < 1 or box["height"] < 1:
                        continue
                    key = (round(box["x"]), round(box["y"]))
                    if key in seen_positions:
                        continue
                    seen_positions.add(key)
                    title = (el.get_attribute("title") or "").strip() \
                            or (el.get_attribute("aria-label") or "").strip()
                    tag_info = el.evaluate(
                        "e => e.tagName + (e.className && typeof e.className === 'string' "
                        "? '.' + e.className.replace(/ /g, '.').slice(0, 60) : '')"
                    )
                    candidates.append({
                        "el": el, "title": title, "tag": tag_info, "box": box,
                    })
                except Exception:
                    continue

        if not candidates:
            print(f"[scrape]     ⚠ no download candidates found in category {cat_label}")
            # Dump full HTML for this category so we can inspect manually
            try:
                safe = re.sub(r"[^a-zA-Z0-9]+", "_", cat_label)
                (DUMP_DIR / f"docs_center_{cat_idx:02d}_{safe}.html").write_text(
                    page.content(), encoding="utf-8"
                )
                print(f"[scrape]       (dumped HTML to docs_center_{cat_idx:02d}_{safe}.html)")
            except Exception:
                pass
            continue

        print(f"[scrape]     {len(candidates)} unique download candidates")

        # ---- 2c. Click each candidate, log result ----
        # IMPORTANT: page.expect_download() CONSUMES the download event for
        # its scope, so the context-level `_on_download` listener does NOT
        # fire for downloads triggered inside `with page.expect_download()`.
        # We must explicitly save the download object via dl_info.value
        # inside the with-block — otherwise the click "succeeds" but no file
        # ever lands on disk.
        def _save_download(download_obj, source_label: str) -> bool:
            try:
                suggested = download_obj.suggested_filename or "unnamed_download"
                save_path = _unique_path(downloads_dir / _safe_filename(suggested))
                download_obj.save_as(str(save_path))
                size = save_path.stat().st_size if save_path.exists() else 0
                captured_downloads.append({
                    "saved_path": save_path,
                    "suggested_filename": suggested,
                    "url": download_obj.url,
                    "source_section": source_label,
                    "link_text": "",
                })
                print(f"[scrape]       ⬇ saved {save_path.name} ({size / 1024:.0f} KB)")
                return True
            except Exception as e:
                print(f"[scrape]       ⬇ save FAILED: {type(e).__name__}: {e}")
                return False

        for i, c in enumerate(candidates):
            title_short = c["title"][:50] or "(no title)"
            tag_short = c["tag"][:50] if c["tag"] else "?"
            print(f"[scrape]     [{i + 1}/{len(candidates)}] {tag_short} title='{title_short}'")
            total_click_attempts += 1

            popups_before = len(popup_urls)

            # Strategy 1: standard click wrapped in expect_download
            saved = False
            try:
                with page.expect_download(timeout=5000) as dl_info:
                    c["el"].click(timeout=2000, force=True)
                saved = _save_download(dl_info.value, f"文档中心 / {cat_label}")
                if saved:
                    continue
            except PWTimeout:
                pass
            except Exception as e:
                print(f"[scrape]       strategy 1 error: {type(e).__name__}: {e}")

            # Did the click open a popup instead?
            if len(popup_urls) > popups_before:
                print(f"[scrape]       ⓘ popup, no download fired")
                continue

            # Strategy 2: hover first, then click
            try:
                c["el"].hover(timeout=1000)
                page.wait_for_timeout(300)
                with page.expect_download(timeout=4000) as dl_info:
                    c["el"].click(timeout=2000, force=True)
                if _save_download(dl_info.value, f"文档中心 / {cat_label}"):
                    continue
            except PWTimeout:
                pass
            except Exception:
                pass

            # Strategy 3: JS-triggered click (bypasses any pointer-event issues)
            try:
                with page.expect_download(timeout=4000) as dl_info:
                    c["el"].evaluate("el => el.click()")
                if _save_download(dl_info.value, f"文档中心 / {cat_label}"):
                    continue
            except PWTimeout:
                pass
            except Exception:
                pass

            print(f"[scrape]       ✗ no download fired (3 strategies tried)")

    new_count = len(captured_downloads) - initial_count
    print(f"\n[scrape]   ========== HARVEST DONE ==========")
    print(f"[scrape]   {len(category_locators)} categories walked")
    print(f"[scrape]   {total_click_attempts} click attempts")
    print(f"[scrape]   {new_count} downloads captured")
    if popup_urls:
        print(f"[scrape]   {len(popup_urls)} popups opened (URLs saved for review)")
        (DUMP_DIR / "popup_urls.txt").write_text("\n".join(popup_urls), encoding="utf-8")

    # If auto-mode produced zero downloads, offer manual fallback
    if new_count == 0:
        print()
        print("[scrape]   AUTO HARVEST FOUND ZERO DOWNLOADS.")
        print("[scrape]   The browser is still open. Switch to it and CLICK the")
        print("[scrape]   download icons manually — every download you trigger")
        print("[scrape]   will be auto-saved by the listener.")
        print("[scrape]   When you're done, return here and press Enter.")
        try:
            input("[scrape]   waiting for you to download manually... press Enter when done: ")
        except EOFError:
            pass
        new_count = len(captured_downloads) - initial_count
        print(f"[scrape]   {new_count} downloads captured (auto + manual)")

    return new_count


def _extract_download_urls_from_apis(api_responses: list[dict]) -> list[dict]:
    """Mine the captured XHR JSON for any string fields that look like file URLs."""
    found = []
    seen = set()
    pattern = re.compile(r'"(https?://[^"]+\.(?:pdf|docx?|xlsx?|pptx?|zip|csv))"', re.IGNORECASE)
    for resp in api_responses:
        body = resp.get("body", "") or ""
        for match in pattern.findall(body):
            if match not in seen:
                seen.add(match)
                found.append({"url": match, "text": "(from API)"})
    return found


# --- Download a file via the authenticated session -------------------------

def _filename_from_disposition(cd: str) -> str | None:
    if not cd:
        return None
    m = re.search(r"filename\*=UTF-8''([^;]+)", cd, re.IGNORECASE)
    if m:
        return unquote(m.group(1).strip().strip('"'))
    m = re.search(r'filename="?([^";]+)"?', cd, re.IGNORECASE)
    if m:
        return m.group(1).strip()
    return None


def _filename_from_url(url: str) -> str | None:
    path = urlparse(url).path
    name = unquote(path.rsplit("/", 1)[-1])
    return name or None


def _safe_filename(name: str) -> str:
    name = re.sub(r'[<>:"/\\|?*\x00-\x1f]', '_', name).strip()
    return name[:200] or "unnamed"


def _unique_path(path: Path) -> Path:
    if not path.exists():
        return path
    stem, suffix = path.stem, path.suffix
    n = 1
    while True:
        cand = path.with_name(f"{stem}__{n}{suffix}")
        if not cand.exists():
            return cand
        n += 1


def _download_file(context, url: str, save_dir: Path, hint_text: str = "") -> Path | None:
    try:
        response = context.request.get(url, timeout=60000)
    except Exception as e:
        print(f"[scrape]     download error ({type(e).__name__}): {url[:80]}")
        return None
    if response.status >= 400:
        print(f"[scrape]     {response.status}: {url[:80]}")
        return None
    try:
        body = response.body()
    except Exception as e:
        print(f"[scrape]     body read failed: {e}")
        return None
    if len(body) > MAX_FILE_BYTES:
        print(f"[scrape]     skipping ({len(body) / 1e6:.1f}MB > {MAX_FILE_BYTES / 1e6:.0f}MB cap)")
        return None

    cd = response.headers.get("content-disposition", "")
    filename = (
        _filename_from_disposition(cd)
        or _filename_from_url(url)
        or (_safe_filename(hint_text) if hint_text else None)
        or "unnamed_download"
    )
    filename = _safe_filename(filename)
    save_path = _unique_path(save_dir / filename)
    save_path.write_bytes(body)
    print(f"[scrape]     saved {save_path.name} ({len(body) / 1024:.0f} KB)")
    return save_path


# --- Document text extraction -----------------------------------------------

def _extract_text(file_path: Path) -> str:
    """Dispatch to the right extractor based on extension. Returns extracted
    text (capped to PER_DOC_TEXT_CHARS) or an error sentinel."""
    ext = file_path.suffix.lower()
    try:
        if ext == ".pdf":
            return _extract_pdf(file_path)
        if ext == ".docx":
            return _extract_docx(file_path)
        if ext in (".xlsx", ".xlsm"):
            return _extract_xlsx(file_path)
        if ext == ".pptx":
            return _extract_pptx(file_path)
        if ext in (".txt", ".md", ".csv"):
            return file_path.read_text(encoding="utf-8", errors="ignore")[:PER_DOC_TEXT_CHARS]
    except Exception as e:
        return f"[extraction failed: {type(e).__name__}: {e}]"
    return f"[unsupported file type: {ext}]"


def _extract_pdf(path: Path) -> str:
    if pdfplumber is None:
        return "[pdfplumber not installed — run: pip install pdfplumber]"
    parts = []
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages):
            t = page.extract_text() or ""
            if t.strip():
                parts.append(f"-- Page {i + 1} --\n{t.strip()}")
            if sum(len(p) for p in parts) > PER_DOC_TEXT_CHARS:
                break
    return "\n\n".join(parts)[:PER_DOC_TEXT_CHARS]


def _extract_docx(path: Path) -> str:
    if DocxDocument is None:
        return "[python-docx not installed — run: pip install python-docx]"
    doc = DocxDocument(path)
    parts = []
    for p in doc.paragraphs:
        t = p.text.strip()
        if t:
            parts.append(t)
    for tbl in doc.tables:
        for row in tbl.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
        if sum(len(p) for p in parts) > PER_DOC_TEXT_CHARS:
            break
    return "\n\n".join(parts)[:PER_DOC_TEXT_CHARS]


def _extract_xlsx(path: Path) -> str:
    if load_workbook is None:
        return "[openpyxl not installed — run: pip install openpyxl]"
    wb = load_workbook(path, data_only=True, read_only=True)
    parts = []
    for sn in wb.sheetnames:
        parts.append(f"### Sheet: {sn}")
        ws = wb[sn]
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            cells = ["" if v is None else str(v).strip() for v in row]
            if any(cells):
                parts.append(" | ".join(cells))
            if i > 200:
                parts.append("[truncated at 200 rows]")
                break
        if sum(len(p) for p in parts) > PER_DOC_TEXT_CHARS:
            break
    return "\n".join(parts)[:PER_DOC_TEXT_CHARS]


def _extract_pptx(path: Path) -> str:
    if PptxPresentation is None:
        return "[python-pptx not installed — run: pip install python-pptx]"
    prs = PptxPresentation(path)
    parts = []
    for i, slide in enumerate(prs.slides):
        slide_lines = [f"-- Slide {i + 1} --"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(r.text for r in para.runs).strip()
                    if t:
                        slide_lines.append(t)
        if len(slide_lines) > 1:
            parts.append("\n".join(slide_lines))
        if sum(len(p) for p in parts) > PER_DOC_TEXT_CHARS:
            break
    return "\n\n".join(parts)[:PER_DOC_TEXT_CHARS]


# --- Markdown distillation --------------------------------------------------

_CHROME_LINES = {
    "版本门户", "ENTERPRISE SYSTEMS", "主导航", "支持与学习",
    "管理员用户", "admin", "系统自检", "退出登录",
}


def _strip_chrome(text: str) -> str:
    lines = text.split("\n")
    cleaned = [ln for ln in lines if ln.strip() not in _CHROME_LINES]
    out = []
    blank = False
    for ln in cleaned:
        if not ln.strip():
            if not blank:
                out.append("")
            blank = True
        else:
            out.append(ln)
            blank = False
    return "\n".join(out).strip()


def _distill_pages(sections: list[dict], scraped_at: str) -> str:
    lines = [
        "# Streamax Internal Version Portal — Page Content",
        "",
        f"_Scraped from `{BASE_URL}` at `{scraped_at}` by `scrape_roadmap.py`._",
        "_Source is internal Streamax intranet — content reflects what an authenticated admin sees._",
        "",
        "Contents are auto-generated from the 8 sidebar sections of the internal version portal:",
        "active roadmap, version history, version logs, feature list, docs index, feature descriptions,",
        "MDVR operation config, and common links. When asked about upcoming releases, what changed in",
        "a version, what features a product supports, or where to find documentation — ground answers",
        "in this material. Document text content is in the companion file `09_roadmap_documents.md`.",
        "",
        "---",
        "",
    ]
    for s in sections:
        if not s.get("ok"):
            lines.append(f"## {s['label']}")
            lines.append(f"_Scrape error: {s.get('error', 'unknown')}_")
            lines.append("")
            continue
        url = s.get("captured_url", BASE_URL)
        text = _strip_chrome(s.get("visible_text", "").strip())
        lines.append(f"## {s['label']}")
        lines.append(f"_Source path: `{url.replace(BASE_URL, '')}`_")
        lines.append("")
        if text:
            lines.append("```")
            lines.append(text[:25000])
            if len(text) > 25000:
                lines.append("\n[... truncated for length ...]")
            lines.append("```")
        else:
            lines.append("_(empty page — possibly requires interaction or had a render issue)_")
        lines.append("")
    return "\n".join(lines)


def _distill_documents(downloads: list[dict], scraped_at: str) -> str:
    """Build the document knowledge file. Each entry: filename + source section
    + extracted text excerpt."""
    lines = [
        "# Streamax Internal Version Portal — Downloaded Documents",
        "",
        f"_Scraped from `{BASE_URL}` at `{scraped_at}` by `scrape_roadmap.py`._",
        "",
        f"This file contains text extracted from {len(downloads)} documents downloaded from",
        "the internal version portal. PDFs, Word docs, Excel sheets, and PowerPoint decks are",
        "all parsed; per-document text is capped at "
        f"~{PER_DOC_TEXT_CHARS:,} characters. When a user asks about a specific spec, manual,",
        "release note, or any other document on the portal, ground the answer in this file's",
        "extracts. Reference documents by filename when relevant.",
        "",
        "---",
        "",
        "## Document index",
        "",
        "| File | Source section | Type | Source page |",
        "| --- | --- | --- | --- |",
    ]
    for d in downloads:
        path = d["saved_path"]
        size_kb = path.stat().st_size // 1024 if path.exists() else 0
        lines.append(
            f"| `{path.name}` | {d['source_section']} | {path.suffix.lstrip('.').upper()} | "
            f"`{d.get('source_url', '')[-60:]}` ({size_kb:,} KB) |"
        )
    lines.append("")
    lines.append("---")
    lines.append("")

    for d in downloads:
        path = d["saved_path"]
        lines.append(f"## File: `{path.name}`")
        lines.append(f"_From section: **{d['source_section']}**_")
        if d.get("link_text"):
            lines.append(f"_Link text in portal: \"{d['link_text']}\"_")
        if d.get("source_url"):
            lines.append(f"_Source URL: `{d['source_url']}`_")
        lines.append("")
        text = d.get("extracted_text", "").strip()
        if text:
            lines.append("```")
            lines.append(text)
            lines.append("```")
        else:
            lines.append("_(no text extracted)_")
        lines.append("")
        lines.append("---")
        lines.append("")
    return "\n".join(lines)


# --- Main -------------------------------------------------------------------

def main(headless: bool = False) -> int:
    DUMP_DIR.mkdir(exist_ok=True)
    DOWNLOADS_DIR.mkdir(exist_ok=True)
    KNOWLEDGE_DIR.mkdir(exist_ok=True)

    scraped_at = datetime.now(timezone.utc).isoformat(timespec="seconds")
    captured_apis: list[dict] = []

    # Downloads triggered by ANY click during the session land here.
    # Populated by the context-level `download` listener below.
    captured_downloads: list[dict] = []

    with sync_playwright() as p:
        browser = p.chromium.launch(headless=headless, slow_mo=50 if not headless else 0)
        context = browser.new_context(
            viewport={"width": 1600, "height": 1000},
            ignore_https_errors=True,
            accept_downloads=True,
        )
        page = context.new_page()

        def _on_response(response):
            try:
                ct = response.headers.get("content-type", "")
                url = response.url
                if "json" in ct or "/api/" in url or url.endswith(".json"):
                    body = response.text()
                    captured_apis.append({
                        "url": url, "status": response.status,
                        "content_type": ct, "body": body[:80000],
                    })
            except Exception:
                pass
        page.on("response", _on_response)

        def _on_download(download):
            """Auto-save any download triggered anywhere in the session."""
            suggested = download.suggested_filename or "unnamed_download"
            save_path = _unique_path(DOWNLOADS_DIR / _safe_filename(suggested))
            try:
                download.save_as(str(save_path))
                size = save_path.stat().st_size if save_path.exists() else 0
                captured_downloads.append({
                    "saved_path": save_path,
                    "suggested_filename": suggested,
                    "url": download.url,
                    "source_section": "(captured live)",
                    "link_text": "",
                })
                print(f"[scrape]     ⬇ saved {save_path.name} ({size / 1024:.0f} KB)")
            except Exception as e:
                print(f"[scrape]     ⬇ FAILED to save {suggested}: {e}")
        context.on("download", _on_download)

        # 1. LOGIN
        if not _try_login(page):
            print("\n[scrape] auto-login didn't find the form — please log in manually")
            print(f"          in the browser window with: {LOGIN_USER} / {LOGIN_PASS}")
            input("          press Enter once you're past the login screen ... ")
        try:
            page.screenshot(path=str(DUMP_DIR / "00_after_login.png"), full_page=True)
        except Exception:
            pass

        # 2. CRAWL SIDEBAR
        sections: list[dict] = []
        for slug, label, _desc in SIDEBAR_ITEMS:
            sections.append(_capture_section(page, slug, label, DUMP_DIR))

            # When we land on 文档中心, do the deeper download-harvest pass —
            # click each category, expand each model row, click each cloud
            # icon. The context-level download listener auto-saves the files.
            if label == "文档中心" and sections[-1].get("ok"):
                try:
                    _harvest_docs_center(page, DOWNLOADS_DIR, captured_downloads)
                except Exception as e:
                    print(f"[scrape]   docs-center harvest failed: {e}")

        # 3. ALSO TRY HREF-BASED DOWNLOADS (cheap fallback for any pages
        #    that DO use plain <a href>, e.g. 常用链接)
        all_links: list[dict] = []
        seen_urls = set()
        for s in sections:
            for link in s.get("download_links", []) or []:
                if link["url"] not in seen_urls:
                    seen_urls.add(link["url"])
                    all_links.append({
                        "url": link["url"],
                        "text": link.get("text", ""),
                        "source_section": s["label"],
                    })
        for link in _extract_download_urls_from_apis(captured_apis):
            if link["url"] not in seen_urls:
                seen_urls.add(link["url"])
                all_links.append({
                    "url": link["url"], "text": link.get("text", ""),
                    "source_section": "(discovered via API)",
                })

        if all_links:
            print(f"\n[scrape] also fetching {len(all_links)} href-based downloads")
            for i, link in enumerate(all_links[:MAX_TOTAL_FILES]):
                print(f"[scrape] [{i + 1}/{min(len(all_links), MAX_TOTAL_FILES)}] "
                      f"{link['source_section']}: {link['url'][:80]}")
                saved = _download_file(context, link["url"], DOWNLOADS_DIR, link.get("text", ""))
                if saved:
                    captured_downloads.append({
                        "saved_path": saved,
                        "source_url": link["url"],
                        "link_text": link.get("text", ""),
                        "source_section": link["source_section"],
                    })

        # 4. DEDUPE downloads by saved path (live + href passes can overlap
        #    if the portal happens to expose both)
        seen_paths = set()
        downloads: list[dict] = []
        for d in captured_downloads:
            p = str(d["saved_path"])
            if p not in seen_paths:
                seen_paths.add(p)
                downloads.append(d)

        browser.close()

    # 5. DUMP API CAPTURES
    api_path = DUMP_DIR / "captured_apis.json"
    api_path.write_text(json.dumps(captured_apis, indent=2, ensure_ascii=False), encoding="utf-8")
    print(f"\n[scrape] {len(captured_apis)} API responses → {api_path.relative_to(ROOT)}")

    # 6. EXTRACT DOC TEXT (no browser needed for this)
    print(f"\n[scrape] extracting text from {len(downloads)} downloaded files...")
    for d in downloads:
        d["extracted_text"] = _extract_text(d["saved_path"])

    # 7. WRITE BOTH KNOWLEDGE FILES
    pages_md = _distill_pages(sections, scraped_at)
    PAGES_OUTPUT_FILE.write_text(pages_md, encoding="utf-8")
    print(f"[scrape] page content ({len(pages_md):,} chars) → {PAGES_OUTPUT_FILE.relative_to(ROOT)}")

    docs_md = _distill_documents(downloads, scraped_at)
    DOCS_OUTPUT_FILE.write_text(docs_md, encoding="utf-8")
    print(f"[scrape] document content ({len(docs_md):,} chars) → {DOCS_OUTPUT_FILE.relative_to(ROOT)}")

    # 8. NEXT STEPS
    ok_sections = sum(1 for s in sections if s.get("ok"))
    print()
    print(f"[scrape] {ok_sections}/{len(sections)} sections captured")
    print(f"[scrape] {len(downloads)} files downloaded + extracted "
          f"({len(all_links)} via href fallback, rest via live cloud-icon clicks)")
    print()
    print("Next steps:")
    print(f"  1. Inspect _scrape_dump/ — screenshots, raw HTML, captured_apis.json,")
    print(f"     downloads/ contains every file pulled from the portal")
    print(f"  2. Review the two distilled files:")
    print(f"       jerry_gpt_knowledge/08_roadmap_portal.md")
    print(f"       jerry_gpt_knowledge/09_roadmap_documents.md")
    print(f"  3. Commit + push:")
    print(f"       git add jerry_gpt_knowledge/08_roadmap_portal.md "
          f"jerry_gpt_knowledge/09_roadmap_documents.md")
    print(f"       git commit -m \"Refresh Jerry roadmap-portal knowledge ({scraped_at[:10]})\"")
    print(f"       git push origin main")
    return 0


if __name__ == "__main__":
    headless = "--headless" in sys.argv
    sys.exit(main(headless=headless))
