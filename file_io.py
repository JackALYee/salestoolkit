"""File I/O for Jerry GPT — uploads in, documents out.

INPUT  (build_user_content): turn Streamlit-uploaded files into Anthropic
content blocks. Images → image blocks, PDFs → document blocks (Claude reads
PDFs natively), Office files (docx/xlsx/pptx) → extracted text, plain text →
decoded text. Pattern adapted from the PM workstation's msg_to_anthropic().

OUTPUT (extract_artifacts / render_artifact): Jerry emits a fenced
```artifact ...``` block of JSON when the user asks for a downloadable
document; we parse it and render a real .docx / .pptx / .xlsx / .pdf with
python-docx / python-pptx / openpyxl / reportlab. strip_artifacts() removes
the JSON block from the displayed text so the user sees prose + a download
button, not raw JSON.

Everything is best-effort and defensive — a bad file or malformed artifact
returns a readable error string, never raises into the chat loop.
"""
from __future__ import annotations

import base64
import io
import json
import re

_MAX_EXTRACT_CHARS = 60000
_IMAGE_MIMES = {"image/png", "image/jpeg", "image/gif", "image/webp"}


# ===========================================================================
# INPUT — uploaded files → Anthropic content blocks
# ===========================================================================

def _ext(name: str) -> str:
    return name.rsplit(".", 1)[-1].lower() if "." in name else ""


def _extract_docx(raw: bytes) -> str:
    try:
        import docx
        d = docx.Document(io.BytesIO(raw))
        parts = [p.text for p in d.paragraphs if p.text.strip()]
        for tbl in d.tables:
            for row in tbl.rows:
                cells = [c.text.strip() for c in row.cells]
                if any(cells):
                    parts.append(" | ".join(cells))
        return "\n".join(parts)
    except Exception as e:
        return f"(could not parse Word document: {e})"


def _extract_xlsx(raw: bytes) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(raw), data_only=True, read_only=True)
        out = []
        for ws in wb.worksheets:
            out.append(f"## Sheet: {ws.title}")
            for row in ws.iter_rows(values_only=True):
                vals = ["" if v is None else str(v) for v in row]
                if any(v.strip() for v in vals):
                    out.append(" | ".join(vals))
        return "\n".join(out)
    except Exception as e:
        return f"(could not parse Excel workbook: {e})"


def _extract_pptx(raw: bytes) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(raw))
        out = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    texts.append(shape.text_frame.text.strip())
            if texts:
                out.append(f"## Slide {i}\n" + "\n".join(texts))
        return "\n\n".join(out)
    except Exception as e:
        return f"(could not parse PowerPoint deck: {e})"


def file_to_block(name: str, mime: str, raw: bytes) -> dict:
    """Convert one uploaded file's bytes into an Anthropic content block."""
    ext = _ext(name)
    # Images → native image block
    if mime in _IMAGE_MIMES or ext in ("png", "jpg", "jpeg", "gif", "webp"):
        media = mime if mime in _IMAGE_MIMES else f"image/{'jpeg' if ext in ('jpg','jpeg') else ext}"
        return {
            "type": "image",
            "source": {"type": "base64", "media_type": media,
                       "data": base64.b64encode(raw).decode("ascii")},
        }
    # PDFs → native document block (Claude reads PDFs directly)
    if mime == "application/pdf" or ext == "pdf":
        return {
            "type": "document",
            "source": {"type": "base64", "media_type": "application/pdf",
                       "data": base64.b64encode(raw).decode("ascii")},
        }
    # Office files → extracted text
    if ext == "docx":
        txt = _extract_docx(raw)
    elif ext in ("xlsx", "xlsm"):
        txt = _extract_xlsx(raw)
    elif ext == "pptx":
        txt = _extract_pptx(raw)
    else:
        try:
            txt = raw.decode("utf-8", "ignore")
        except Exception:
            txt = f"(unsupported file type: {name})"
    return {"type": "text", "text": f"【Attached file: {name}】\n{txt[:_MAX_EXTRACT_CHARS]}"}


def build_user_content(text: str, files: list) -> tuple:
    """Return (api_content, display_note).

    `files` is a list of Streamlit UploadedFile objects (have .name, .type,
    .getvalue()). api_content is a plain string when there are no files, else a
    list of content blocks (text first, then one block per file). display_note
    is a short markdown line summarizing attachments for the UI / history.
    """
    if not files:
        return (text or ""), ""
    blocks = []
    if text:
        blocks.append({"type": "text", "text": text})
    names = []
    for f in files:
        try:
            raw = f.getvalue()
            blocks.append(file_to_block(f.name, getattr(f, "type", "") or "", raw))
            names.append(f.name)
        except Exception as e:
            blocks.append({"type": "text", "text": f"【Attachment {getattr(f,'name','file')} failed: {e}】"})
    note = "📎 *Attached: " + ", ".join(names) + "*" if names else ""
    return blocks, note


# ===========================================================================
# OUTPUT — Jerry's ```artifact``` JSON → real downloadable files
# ===========================================================================

_ARTIFACT_RE = re.compile(r"```artifact\s*\n(.*?)```", re.DOTALL)

_MIME = {
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "pdf": "application/pdf",
}


def extract_artifacts(text: str) -> list[dict]:
    """Find and parse every ```artifact``` JSON block in `text`. Skips blocks
    that don't parse or lack a supported format."""
    if not text or "```artifact" not in text:
        return []
    out = []
    for m in _ARTIFACT_RE.finditer(text):
        try:
            spec = json.loads(m.group(1).strip())
        except Exception:
            continue
        if isinstance(spec, dict) and spec.get("format") in _MIME:
            out.append(spec)
    return out


def strip_artifacts(text: str) -> str:
    """Remove ```artifact``` blocks from text so the user sees prose, not JSON.
    Leaves a small marker where each block was."""
    if not text or "```artifact" not in text:
        return text
    return _ARTIFACT_RE.sub("\n*[document prepared — see the download button below]*\n", text).strip()


def _safe_filename(spec: dict) -> str:
    base = str(spec.get("filename") or spec.get("title") or "Streamax-document").strip()
    base = re.sub(r"[^\w\-. ]+", "", base).strip() or "Streamax-document"
    return f"{base}.{spec['format']}"


def render_artifact(spec: dict) -> tuple:
    """Render one artifact spec → (bytes, filename, mime). Raises on failure
    (caller wraps in try/except)."""
    fmt = spec["format"]
    if fmt == "docx":
        data = _render_docx(spec)
    elif fmt == "pptx":
        data = _render_pptx(spec)
    elif fmt == "xlsx":
        data = _render_xlsx(spec)
    elif fmt == "pdf":
        data = _render_pdf(spec)
    else:
        raise ValueError(f"unsupported format {fmt}")
    return data, _safe_filename(spec), _MIME[fmt]


# --- DOCX ------------------------------------------------------------------
def _render_docx(spec: dict) -> bytes:
    import docx
    d = docx.Document()
    if spec.get("title"):
        d.add_heading(str(spec["title"]), level=0)
    for blk in spec.get("blocks", []):
        t = blk.get("type")
        if t == "heading":
            d.add_heading(str(blk.get("text", "")), level=int(blk.get("level", 1)))
        elif t == "paragraph":
            d.add_paragraph(str(blk.get("text", "")))
        elif t == "bullets":
            for it in blk.get("items", []):
                d.add_paragraph(str(it), style="List Bullet")
        elif t == "table":
            headers = blk.get("headers", [])
            rows = blk.get("rows", [])
            ncol = len(headers) or (len(rows[0]) if rows else 0)
            if ncol:
                tbl = d.add_table(rows=0, cols=ncol)
                tbl.style = "Light Grid Accent 1"
                if headers:
                    cells = tbl.add_row().cells
                    for i, h in enumerate(headers[:ncol]):
                        cells[i].text = str(h)
                for r in rows:
                    cells = tbl.add_row().cells
                    for i, v in enumerate(list(r)[:ncol]):
                        cells[i].text = "" if v is None else str(v)
    buf = io.BytesIO()
    d.save(buf)
    return buf.getvalue()


# --- PPTX ------------------------------------------------------------------
def _render_pptx(spec: dict) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    prs = Presentation()
    # Title slide
    if spec.get("title"):
        s = prs.slides.add_slide(prs.slide_layouts[0])
        s.shapes.title.text = str(spec["title"])
        if spec.get("subtitle") and len(s.placeholders) > 1:
            s.placeholders[1].text = str(spec["subtitle"])
    for slide in spec.get("slides", []):
        layout = prs.slide_layouts[1]  # title + content
        s = prs.slides.add_slide(layout)
        s.shapes.title.text = str(slide.get("title", ""))
        bullets = slide.get("bullets", [])
        body = slide.get("body", "")
        if bullets and len(s.placeholders) > 1:
            tf = s.placeholders[1].text_frame
            tf.clear()
            for i, b in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = str(b)
                p.font.size = Pt(18)
        elif body and len(s.placeholders) > 1:
            s.placeholders[1].text_frame.text = str(body)
        notes = slide.get("notes")
        if notes:
            s.notes_slide.notes_text_frame.text = str(notes)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# --- XLSX ------------------------------------------------------------------
def _render_xlsx(spec: dict) -> bytes:
    import openpyxl
    from openpyxl.styles import Font
    wb = openpyxl.Workbook()
    wb.remove(wb.active)
    sheets = spec.get("sheets") or [{"name": "Sheet1",
                                     "headers": spec.get("headers", []),
                                     "rows": spec.get("rows", [])}]
    for sh in sheets:
        ws = wb.create_sheet(title=str(sh.get("name", "Sheet"))[:31] or "Sheet")
        headers = sh.get("headers", [])
        if headers:
            ws.append([str(h) for h in headers])
            for c in ws[1]:
                c.font = Font(bold=True)
        for r in sh.get("rows", []):
            ws.append(["" if v is None else v for v in r])
    if not wb.worksheets:
        wb.create_sheet(title="Sheet1")
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# --- PDF (reportlab + built-in CJK CID font, so Chinese renders) -----------
def _render_pdf(spec: dict) -> bytes:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import mm
    from reportlab.lib import colors
    from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer,
                                    ListFlowable, ListItem, Table, TableStyle)
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.cidfonts import UnicodeCIDFont

    # Register a CID font that covers Simplified Chinese + Latin. Built into
    # reportlab — no TTF to bundle. Fall back to Helvetica if unavailable.
    font = "Helvetica"
    try:
        pdfmetrics.registerFont(UnicodeCIDFont("STSong-Light"))
        font = "STSong-Light"
    except Exception:
        pass

    styles = getSampleStyleSheet()
    body = ParagraphStyle("body", parent=styles["Normal"], fontName=font,
                          fontSize=11, leading=16)
    h0 = ParagraphStyle("h0", parent=styles["Title"], fontName=font, fontSize=20, leading=24)
    h = lambda lvl: ParagraphStyle(f"h{lvl}", parent=styles["Heading1"], fontName=font,
                                   fontSize=max(12, 18 - lvl), leading=20, spaceBefore=10)

    flow = []
    if spec.get("title"):
        flow += [Paragraph(str(spec["title"]), h0), Spacer(1, 8)]
    for blk in spec.get("blocks", []):
        t = blk.get("type")
        if t == "heading":
            flow.append(Paragraph(str(blk.get("text", "")), h(int(blk.get("level", 1)))))
        elif t == "paragraph":
            flow += [Paragraph(str(blk.get("text", "")), body), Spacer(1, 6)]
        elif t == "bullets":
            items = [ListItem(Paragraph(str(it), body)) for it in blk.get("items", [])]
            flow += [ListFlowable(items, bulletType="bullet"), Spacer(1, 6)]
        elif t == "table":
            headers = blk.get("headers", [])
            rows = blk.get("rows", [])
            data = ([[str(x) for x in headers]] if headers else []) + \
                   [["" if v is None else str(v) for v in r] for r in rows]
            if data:
                tbl = Table(data, hAlign="LEFT")
                tbl.setStyle(TableStyle([
                    ("FONTNAME", (0, 0), (-1, -1), font),
                    ("FONTSIZE", (0, 0), (-1, -1), 9),
                    ("GRID", (0, 0), (-1, -1), 0.4, colors.grey),
                    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#E6F0EA")),
                    ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ]))
                flow += [tbl, Spacer(1, 8)]
    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4, topMargin=18 * mm, bottomMargin=18 * mm,
                            leftMargin=18 * mm, rightMargin=18 * mm)
    doc.build(flow or [Paragraph(str(spec.get("title", "Document")), body)])
    return buf.getvalue()


# Injected into Jerry's system prompt so he knows how to emit a downloadable doc.
ARTIFACT_HINT = (
    "FILE GENERATION: You can produce downloadable Office/PDF documents. When "
    "the user asks you to make a PPT/deck, Word doc, Excel sheet, or PDF, do "
    "two things: (1) give a short normal reply, and (2) emit ONE fenced code "
    "block tagged `artifact` containing a JSON spec. The interface renders it "
    "into a real file with a download button — never paste a link or base64.\n"
    "Schemas (pick by `format`):\n"
    "- Word/PDF: {\"format\":\"docx\"|\"pdf\",\"filename\":\"...\",\"title\":\"...\","
    "\"blocks\":[{\"type\":\"heading\",\"level\":1,\"text\":\"...\"},"
    "{\"type\":\"paragraph\",\"text\":\"...\"},"
    "{\"type\":\"bullets\",\"items\":[\"...\"]},"
    "{\"type\":\"table\",\"headers\":[\"...\"],\"rows\":[[\"...\"]]}]}\n"
    "- PowerPoint: {\"format\":\"pptx\",\"filename\":\"...\",\"title\":\"...\","
    "\"subtitle\":\"...\",\"slides\":[{\"title\":\"...\",\"bullets\":[\"...\"],"
    "\"notes\":\"...\"}]}\n"
    "- Excel: {\"format\":\"xlsx\",\"filename\":\"...\",\"sheets\":[{\"name\":\"...\","
    "\"headers\":[\"...\"],\"rows\":[[\"...\"]]}]}\n"
    "Write real, complete content into the spec (not placeholders). Only emit an "
    "artifact when the user actually wants a file; otherwise answer normally."
)
